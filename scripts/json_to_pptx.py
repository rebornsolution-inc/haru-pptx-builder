#!/usr/bin/env python3
"""
JSON to PPTX Converter - 편집 가능한 PowerPoint 생성

presentation.json 데이터를 기반으로 실제 텍스트/도형을 포함한 PPTX 생성

## 조정 가능한 설정 (% 기반)

제작 시마다 아래 설정을 조정하여 결과물을 미세 조정할 수 있습니다:

| 설정 | 기본값 | 설명 |
|------|--------|------|
| FONT_SCALE | 0.95 (-5%) | 폰트 크기 조정 비율 |
| LINE_SPACING_SCALE | 0.83 (-17%) | 줄간격 조정 비율 (기준 1.2 대비) |
| PARAGRAPH_SPACING_SCALE | 0.0 | 문단 간격 조정 비율 (기준 폰트 대비) |
| IMAGE_CORNER_RATIO | 0.05 (5%) | 이미지 라운딩 비율 |

## 사용법

    python json_to_pptx.py <presentation.json> [output.pptx]

## 예시

    python json_to_pptx.py projects/eumlogistic/presentation.json
    python json_to_pptx.py projects/eumlogistic/presentation.json output.pptx
"""

import json
import sys
import os
import re
import requests
import shutil
from io import BytesIO
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

from PIL import Image, ImageDraw, ImageFont


# ========================================
# 📐 조정 가능한 설정 (% 기반)
# ========================================
# 제작 시마다 이 값들을 조정하여 결과물을 미세 조정할 수 있습니다.
# 1.0 = 100% (원본), 0.95 = 95% (-5%), 1.1 = 110% (+10%)

SLIDE_WIDTH = Inches(13.333)   # 16:9 비율 (고정)
SLIDE_HEIGHT = Inches(7.5)

# [폰트 크기] 1.0 = 원본, 0.95 = 5% 축소, 1.1 = 10% 확대
FONT_SCALE = 0.95

# [줄간격] 기준값 1.2 대비 비율. 1.0 = 1.2 유지, 0.83 = 1.0으로 축소
LINE_SPACING_SCALE = 0.83  # 1.2 * 0.83 ≈ 1.0

# [문단 간격] 폰트 크기 대비 비율. 0.0 = 없음, 0.5 = 폰트의 50%
PARAGRAPH_SPACING_SCALE = 0.0

# [이미지 라운딩] 이미지 최소변 대비 비율. 0.05 = 5%
IMAGE_CORNER_RATIO = 0.05

# ========================================
# 디자인 토큰 (presentation.json에서 추출)
# ========================================
COLORS = {
    "primary": "5B6CF9",
    "primaryLight": "7B8BFF",
    "primaryDark": "4A5AE8",
    "secondary": "0F172A",
    "secondaryLight": "1E293B",
    "accent": "06B6D4",
    "accentLight": "22D3EE",
    "background": "FFFFFF",
    "backgroundDark": "0F172A",
    "backgroundAlt": "F8FAFC",
    "text": "1E293B",
    "textLight": "64748B",
    "textMuted": "94A3B8",
    "textOnDark": "FFFFFF",
    "highlight": "F59E0B",
    "success": "10B981",
}

FONTS = {
    "heading": "Pretendard",
    "body": "Pretendard",
}

# Pretendard 폰트 파일 경로
FONT_DIR = Path(__file__).parent.parent / "fonts"
FONT_FILES = {
    "regular": FONT_DIR / "Pretendard-Regular.ttf",
    "medium": FONT_DIR / "Pretendard-Medium.ttf",
    "semibold": FONT_DIR / "Pretendard-SemiBold.ttf",
    "bold": FONT_DIR / "Pretendard-Bold.ttf",
}

def check_fonts():
    """Pretendard 폰트 파일 존재 여부 확인"""
    missing = []
    for weight, path in FONT_FILES.items():
        if not path.exists():
            missing.append(f"{weight}: {path}")
    
    if missing:
        print("\n⚠️  Pretendard 폰트 파일이 없습니다!")
        print("다음 파일들을 fonts/ 폴더에 추가해주세요:\n")
        for m in missing:
            print(f"  - {m}")
        print(f"\n설치 방법은 fonts/README.md를 참고하세요.")
        print("폴백 폰트 '맑은 고딕'을 사용합니다.\n")
        return False
    return True

# 폰트 체크 (스크립트 로드 시 확인)
PRETENDARD_AVAILABLE = check_fonts()


def copy_fonts_to_project(project_dir: Path):
    """사용된 폰트 파일을 프로젝트 폴더에 복사
    
    Args:
        project_dir: 프로젝트 폴더 경로 (예: projects/eumlogistic/)
    """
    if not PRETENDARD_AVAILABLE:
        return
    
    # 프로젝트 내 fonts 폴더 생성
    target_font_dir = project_dir / "fonts"
    target_font_dir.mkdir(exist_ok=True)
    
    # 사용된 폰트 파일 복사
    copied = []
    for weight, source_path in FONT_FILES.items():
        if source_path.exists():
            target_path = target_font_dir / source_path.name
            shutil.copy2(source_path, target_path)
            copied.append(source_path.name)
    
    if copied:
        print(f"\n[Fonts] 프로젝트에 폰트 복사 완료:")
        for font in copied:
            print(f"  → {target_font_dir / font}")


def hex_to_rgb(hex_color: str) -> RGBColor:
    """HEX 색상을 RGBColor로 변환"""
    hex_color = hex_color.lstrip('#')
    r = int(hex_color[0:2], 16)
    g = int(hex_color[2:4], 16)
    b = int(hex_color[4:6], 16)
    return RGBColor(r, g, b)


def download_image(url: str) -> BytesIO | None:
    """URL에서 이미지 다운로드 (재시도 포함)"""
    for attempt in range(2):  # 최대 2회 시도
        try:
            response = requests.get(url, timeout=5)  # 타임아웃 5초로 단축
            if response.status_code == 200:
                return BytesIO(response.content)
        except Exception as e:
            if attempt == 0:
                continue  # 첫 실패 시 재시도
            print(f"  [WARN] 이미지 다운로드 실패 ({attempt+1}회): {url} - {e}")
    return None


def add_rounded_image(slide, image_url, left, top, width, height, radius=None):
    """라운드 처리된 이미지 추가
    
    IMAGE_CORNER_RATIO 비율로 라운드 코너 적용
    """
    img_data = download_image(image_url)
    if not img_data:
        return None
    
    try:
        # 이미지 열기
        img = Image.open(img_data).convert("RGBA")
        
        # 리사이즈 (PPT 크기에 맞게)
        target_width = int(width / Inches(1) * 96)  # 96 DPI
        target_height = int(height / Inches(1) * 96)
        img = img.resize((target_width, target_height), Image.Resampling.LANCZOS)
        
        # 라운드 마스크 생성 - IMAGE_CORNER_RATIO 비율 적용
        mask = Image.new("L", (target_width, target_height), 0)
        draw = ImageDraw.Draw(mask)
        radius_px = int(min(target_width, target_height) * IMAGE_CORNER_RATIO)
        draw.rounded_rectangle([(0, 0), (target_width, target_height)], radius=radius_px, fill=255)
        
        # 투명 배경 생성
        output = Image.new("RGBA", (target_width, target_height), (0, 0, 0, 0))
        output.paste(img, (0, 0), mask)
        
        # BytesIO로 저장
        output_bytes = BytesIO()
        output.save(output_bytes, format="PNG")
        output_bytes.seek(0)
        
        # 슬라이드에 추가
        return slide.shapes.add_picture(output_bytes, left, top, width, height)
    except Exception as e:
        print(f"  [WARN] 이미지 라운딩 처리 실패: {e}")
        # 폴백: 라운딩 없이 추가
        img_data.seek(0)
        return slide.shapes.add_picture(img_data, left, top, width, height)


def add_text_box(slide, left, top, width, height, text, font_size=18, font_color="1E293B", 
                 bold=False, align="left", font_name=None, line_spacing=None):
    """텍스트 박스 추가
    
    FONT_SCALE, LINE_SPACING_SCALE, PARAGRAPH_SPACING_SCALE 적용
    """
    shape = slide.shapes.add_textbox(left, top, width, height)
    tf = shape.text_frame
    tf.word_wrap = True
    
    # 폰트 크기에 스케일 적용
    scaled_font_size = font_size * FONT_SCALE
    
    # 줄간격 기본값: 기준 1.2에 LINE_SPACING_SCALE 적용
    if line_spacing is None:
        line_spacing = 1.2 * LINE_SPACING_SCALE
    
    # 문단 간격: 폰트 크기에 PARAGRAPH_SPACING_SCALE 적용
    paragraph_spacing = scaled_font_size * PARAGRAPH_SPACING_SCALE
    
    # 첫 번째 줄인지 확인
    lines = text.split('\n')
    
    for i, line in enumerate(lines):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        
        p.text = line
        p.font.size = Pt(scaled_font_size)
        p.font.color.rgb = hex_to_rgb(font_color)
        p.font.bold = bold
        
        # Pretendard 폰트 사용 (사용 가능한 경우)
        if PRETENDARD_AVAILABLE:
            p.font.name = "Pretendard"
        else:
            p.font.name = font_name or "맑은 고딕"
        
        p.line_spacing = line_spacing
        p.space_after = Pt(paragraph_spacing)
        p.space_before = Pt(paragraph_spacing)
        
        if align == "center":
            p.alignment = PP_ALIGN.CENTER
        elif align == "right":
            p.alignment = PP_ALIGN.RIGHT
        else:
            p.alignment = PP_ALIGN.LEFT
    
    return shape


def add_rectangle(slide, left, top, width, height, fill_color=None, line_color=None, radius=0):
    """사각형 도형 추가"""
    if radius > 0:
        shape = slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, left, top, width, height)
        # 라운드 조절
        shape.adjustments[0] = min(radius / 100, 0.5)
    else:
        shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, left, top, width, height)
    
    if fill_color:
        shape.fill.solid()
        shape.fill.fore_color.rgb = hex_to_rgb(fill_color)
    else:
        shape.fill.background()
    
    if line_color:
        shape.line.color.rgb = hex_to_rgb(line_color)
    else:
        shape.line.fill.background()
    
    return shape


def set_slide_background(slide, color):
    """슬라이드 배경색 설정"""
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = hex_to_rgb(color)


def set_gradient_background(slide, color1, color2):
    """그라데이션 배경 (단색으로 대체 - python-pptx 제한)"""
    # python-pptx는 복잡한 그라데이션 지원이 제한적이므로 메인 색상 사용
    set_slide_background(slide, color1)


# ===== Helper Functions for Text Extraction =====

def find_text_by_keyword(texts, keyword, text_type=None):
    """텍스트 배열에서 키워드로 텍스트 찾기"""
    for text_obj in texts:
        text = text_obj.get("text", "")
        if keyword in text:
            if text_type is None or text_obj.get("type") == text_type:
                return text
    return ""

def find_texts_by_type(texts, text_type):
    """특정 타입의 모든 텍스트 찾기"""
    return [t.get("text", "") for t in texts if t.get("type") == text_type]

def extract_numbers_from_texts(texts):
    """텍스트에서 숫자 추출 (통계 값 등)"""
    import re
    numbers = []
    for text_obj in texts:
        text = text_obj.get("text", "")
        # 숫자 패턴 찾기: 2020, 35억, 7명, 3.3억 등
        matches = re.findall(r'\d+(?:\.\d+)?(?:억|명|년)?', text)
        numbers.extend(matches)
    return numbers


def create_generic_slide(prs, slide_data, design_tokens):
    """범용 슬라이드 생성 - 간단한 레이아웃"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    
    # 템플릿에 따라 배경색 설정
    template = slide_data.get("template", "")
    if "dark" in template or slide_data.get("slideNumber", 0) in [1, 3, 7, 11]:
        set_slide_background(slide, COLORS["secondary"])
        default_color = "FFFFFF"
    elif slide_data.get("slideNumber", 0) in [4, 6, 9]:
        set_slide_background(slide, COLORS["backgroundAlt"])
        default_color = COLORS["text"]
    elif slide_data.get("slideNumber", 0) == 10:
        set_slide_background(slide, COLORS["primary"])
        default_color = "FFFFFF"
    else:
        set_slide_background(slide, COLORS["background"])
        default_color = COLORS["text"]
    
    elements = slide_data.get("elements", {})
    texts = elements.get("texts", [])
    images = elements.get("images", [])
    
    # 제목 (슬라이드 상단)
    title = slide_data.get("title", "")
    if title:
        add_text_box(slide, Inches(0.8), Inches(0.5), Inches(11), Inches(0.8),
                     title, font_size=36, font_color=default_color, bold=True)
    
    # 본문 텍스트 배치 (단순 나열)
    y_offset = Inches(1.5)
    for i, text_obj in enumerate(texts[:10]):  # 최대 10개
        text = text_obj.get("text", "")
        if text and len(text) > 3:  # 의미 있는 텍스트만
            font_size = 16 if len(text) > 50 else 18
            add_text_box(slide, Inches(0.8), y_offset, Inches(11), Inches(0.6),
                         text[:200], font_size=font_size, font_color=default_color)
            y_offset += Inches(0.7)
            if y_offset > Inches(6.5):  # 슬라이드 하단 제한
                break
    
    # 이미지 배치 (우측 또는 하단)
    if images:
        img_url = images[0].get("src", "")
        if img_url:
            add_rounded_image(slide, img_url, Inches(8), Inches(1.5), Inches(4), Inches(3), radius=16)


# ===== 슬라이드 생성 함수들 =====

def create_slide_1_cover(prs, slide_data, design_tokens):
    """슬라이드 1: 표지 (Hero Cover)"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])  # 빈 슬라이드
    set_slide_background(slide, COLORS["secondaryLight"])
    
    elements = slide_data.get("elements", {})
    texts = elements.get("texts", [])
    images = elements.get("images", [])
    
    # 텍스트 요소에서 추출
    logo_text = "E-UM LOGISTIC"
    badge_text = "✦ 신뢰할 수 있는 물류 파트너"
    main_heading = "글로벌 물류의 새로운 기준"
    company_name = "이음로지스틱"
    description = "해상・항공 포워딩부터 내륙운송, 물류 인프라까지\n고객의 화물에 최적화된 토탈 물류 솔루션을 제공합니다."
    
    # texts 배열에서 텍스트 찾기
    for text_obj in texts:
        text = text_obj.get("text", "")
        text_type = text_obj.get("type", "")
        
        if "E-UM LOGISTIC" in text and text_obj.get("color") == "white":
            logo_text = text
        elif "✦" in text and "신뢰" in text:
            badge_text = text
        elif "글로벌" in text and text_type == "heading":
            main_heading = text
        elif "이음로지스틱" in text and text_type == "subheading":
            company_name = text
        elif "해상" in text and "항공" in text:
            description = text.replace("고객의", "\n고객의")
    
    # 로고
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.6),
                 logo_text, font_size=20, font_color="FFFFFF", bold=True)
    
    # 뱃지
    badge_shape = add_rectangle(slide, Inches(4.5), Inches(2.2), Inches(4.3), Inches(0.5),
                                 fill_color="3D4A66", radius=50)
    add_text_box(slide, Inches(4.5), Inches(2.25), Inches(4.3), Inches(0.5),
                 badge_text, font_size=14, font_color="FFFFFF", align="center")
    
    # 메인 헤딩
    add_text_box(slide, Inches(0.5), Inches(2.9), Inches(12.3), Inches(1),
                 main_heading,
                 font_size=56, font_color="FFFFFF", bold=True, align="center")
    
    # 서브 헤딩 (회사명)
    add_text_box(slide, Inches(0.5), Inches(3.9), Inches(12.3), Inches(0.7),
                 company_name,
                 font_size=44, font_color=COLORS["primaryLight"], bold=True, align="center")
    
    # 설명
    add_text_box(slide, Inches(2.5), Inches(4.8), Inches(8.3), Inches(1),
                 description, font_size=18, font_color=COLORS["textMuted"], align="center")
    
    # 하단 이미지 (에러 발생해도 계속 진행)
    if len(images) >= 1:
        try:
            img1_url = images[0].get("src", "")
            if img1_url:
                add_rounded_image(slide, img1_url, Inches(1), Inches(5.9), Inches(5.5), Inches(1.4), radius=16)
        except Exception as e:
            print(f"  [WARN] 이미지 1 추가 실패: {e}")
    
    if len(images) >= 2:
        try:
            img2_url = images[1].get("src", "")
            if img2_url:
                add_rounded_image(slide, img2_url, Inches(6.8), Inches(5.9), Inches(5.5), Inches(1.4), radius=16)
        except Exception as e:
            print(f"  [WARN] 이미지 2 추가 실패: {e}")


def create_slide_2_company(prs, slide_data, design_tokens):
    """슬라이드 2: 회사개요"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    
    elements = slide_data.get("elements", {})
    texts = elements.get("texts", [])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "COMPANY PROFILE",
                 font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.7),
                 "이음로지스틱 주식회사",
                 font_size=40, font_color=COLORS["text"], bold=True)
    
    # 부제목
    add_text_box(slide, Inches(0.8), Inches(1.65), Inches(10), Inches(0.5),
                 "국제 포워딩 전문 기업으로서 해상・항공 물류의 모든 것을 책임집니다.",
                 font_size=18, font_color=COLORS["textLight"])
    
    # 정보 카드들 (2x2 그리드)
    card_data = [
        ("🏢 본사", "서울 강서구 마곡중앙로 161-8\n두산더랜드파크 B동 510호"),
        ("📍 부산지사", "경남 창원시 진해구 신항7로 63\nMS디스트리파크 용동물류센터 206호"),
        ("📞 연락처", "Tel: +82 2 3662 8150~5\nFax: +82 2 6442 8153"),
        ("✉️ 이메일", "info@e-umlk.com"),
    ]
    
    card_positions = [
        (Inches(0.8), Inches(2.3)),   # 좌상
        (Inches(6.6), Inches(2.3)),   # 우상
        (Inches(0.8), Inches(3.8)),   # 좌하
        (Inches(6.6), Inches(3.8)),   # 우하
    ]
    
    for i, (label, value) in enumerate(card_data):
        left, top = card_positions[i]
        
        # 카드 배경
        add_rectangle(slide, left, top, Inches(5.5), Inches(1.3),
                      fill_color=COLORS["backgroundAlt"], radius=12)
        
        # 라벨
        add_text_box(slide, left + Inches(0.2), top + Inches(0.2), Inches(5), Inches(0.35),
                     label,
                     font_size=13, font_color=COLORS["primary"], bold=True)
        
        # 값
        add_text_box(slide, left + Inches(0.2), top + Inches(0.55), Inches(5.1), Inches(0.7),
                     value, font_size=15, font_color=COLORS["text"])
    
    # 통계 행
    stat_data = [
        ("2020", "설립연도"),
        ("35억", "연평균 매출"),
        ("7명", "전문 인력"),
        ("3.3억", "자본금"),
    ]
    
    stat_left = Inches(1.5)
    stat_top = Inches(5.8)
    stat_width = Inches(2.5)
    
    # 구분선
    add_rectangle(slide, Inches(0.8), Inches(5.6), Inches(11.7), Inches(0.02),
                  fill_color="E2E8F0")
    
    for i, (value, label) in enumerate(stat_data):
        x = stat_left + (i * stat_width)
        add_text_box(slide, x, stat_top, stat_width, Inches(0.6),
                     value,
                     font_size=32, font_color=COLORS["primary"], bold=True, align="center")
        add_text_box(slide, x, stat_top + Inches(0.55), stat_width, Inches(0.4),
                     label,
                     font_size=14, font_color=COLORS["textLight"], align="center")


def create_slide_3_history(prs, slide_data, design_tokens):
    """슬라이드 3: 회사연혁"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["secondary"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.5), Inches(3), Inches(0.4),
                 "HISTORY",
                 font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(10), Inches(0.7),
                 "빠르게 성장하는 물류 파트너",
                 font_size=40, font_color="FFFFFF", bold=True)
    
    # 부제목
    add_text_box(slide, Inches(0.8), Inches(1.6), Inches(10), Inches(0.5),
                 "설립 5년 만에 20개 이상의 주요 고객사와 글로벌 네트워크 구축",
                 font_size=18, font_color=COLORS["textMuted"])
    
    # 타임라인 데이터
    timeline_data = [
        ("2020", "이음로지스틱 설립", "롯데웰푸드 계약\n가온전선 계약\nMS디스트리파크 협력"),
        ("2021", "부산 지사 설립", "신항 물류 거점 확보\n비제이신항 협력"),
        ("2022", "사업 확장", "삼양식품\n통합물류 계약"),
        ("2023", "글로벌 진출", "태국 YC리사이클\n파트너십"),
        ("2024", "인프라 강화", "페트로마인\n물류 인프라 구축"),
        ("2025", "서비스 다각화", "광메탈 운송 계약\n삼성헬스토리 냉동물류"),
    ]
    
    timeline_top = Inches(2.8)
    item_width = Inches(2.0)
    start_left = Inches(0.5)
    
    for i, (year, title, events) in enumerate(timeline_data):
        x = start_left + (i * item_width)
        
        # 점
        add_rectangle(slide, x + Inches(0.85), timeline_top, Inches(0.15), Inches(0.15),
                      fill_color=COLORS["primary"], radius=50)
        
        # 연결선 (마지막 제외)
        if i < len(timeline_data) - 1:
            add_rectangle(slide, x + Inches(1.0), timeline_top + Inches(0.06), 
                         item_width - Inches(0.15), Inches(0.03),
                         fill_color="3D4A66")
        
        # 연도
        add_text_box(slide, x, timeline_top + Inches(0.3), item_width, Inches(0.5),
                     year,
                     font_size=30, font_color=COLORS["primary"], bold=True, align="center")
        
        # 제목
        add_text_box(slide, x, timeline_top + Inches(0.85), item_width, Inches(0.4),
                     title,
                     font_size=16, font_color="FFFFFF", bold=True, align="center")
        
        # 이벤트들
        add_text_box(slide, x, timeline_top + Inches(1.35), item_width, Inches(1.5),
                     events,
                     font_size=13, font_color=COLORS["textMuted"], align="center")


def create_slide_4_values(prs, slide_data, design_tokens):
    """슬라이드 4: 핵심가치"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["backgroundAlt"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.5), Inches(0.5), Inches(12), Inches(0.4),
                 "WHY E-UM",
                 font_size=14, font_color=COLORS["primary"], bold=True, align="center")
    
    # 제목
    add_text_box(slide, Inches(0.5), Inches(0.95), Inches(12.3), Inches(1.2),
                 "신뢰를 잇는 물류,\n가치를 잇는 파트너십",
                 font_size=40, font_color=COLORS["text"], bold=True, align="center")
    
    # 설명
    add_text_box(slide, Inches(2), Inches(2.3), Inches(9.3), Inches(0.7),
                 "NVOCC(Non-Vessel Operating Common Carrier)로서 해상・항공 수출입 화물운송의 풍부한 경험으로 신뢰할 수 있는 서비스를 제공합니다.",
                 font_size=17, font_color=COLORS["textLight"], align="center")
    
    # 가치 카드 데이터
    card_data = [
        ("🛡️", "신뢰성", "20개 이상의 주요 고객사가 선택한 검증된 물류 파트너"),
        ("🌍", "글로벌 네트워크", "5개 대륙, 30개국 이상의 파트너사와 함께하는 글로벌 물류망"),
        ("⚡", "전문성", "해상・항공・특수화물까지 토탈 포워딩 솔루션 제공"),
    ]
    
    card_width = Inches(3.8)
    card_height = Inches(3.5)
    card_gap = Inches(0.4)
    start_left = Inches(0.7)
    card_top = Inches(3.3)
    
    for i, (icon, title, desc) in enumerate(card_data):
        x = start_left + i * (card_width + card_gap)
        
        # 카드 배경
        add_rectangle(slide, x, card_top, card_width, card_height,
                      fill_color="FFFFFF", radius=16)
        
        # 아이콘 배경
        add_rectangle(slide, x + Inches(1.3), card_top + Inches(0.4), 
                      Inches(1.2), Inches(1.2),
                      fill_color="EEF0FE", radius=20)
        add_text_box(slide, x + Inches(1.3), card_top + Inches(0.55), Inches(1.2), Inches(1),
                     icon, font_size=40, align="center")
        
        # 제목
        add_text_box(slide, x + Inches(0.2), card_top + Inches(1.8), card_width - Inches(0.4), Inches(0.5),
                     title,
                     font_size=22, font_color=COLORS["text"], bold=True, align="center")
        
        # 설명
        add_text_box(slide, x + Inches(0.2), card_top + Inches(2.4), card_width - Inches(0.4), Inches(0.9),
                     desc,
                     font_size=15, font_color=COLORS["textLight"], align="center")


def create_slide_5_services(prs, slide_data, design_tokens):
    """슬라이드 5: 서비스소개 (Bento Grid) - 하드코딩"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                 "SERVICES", font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
                 "통합 물류 서비스", font_size=36, font_color=COLORS["text"], bold=True)
    
    # ===== Bento Grid 레이아웃 =====
    
    # 대형 카드: 해상 포워딩 (좌측, 세로로 길게)
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1494412574643-ff11b0a5c1c3?w=400&h=500&fit=crop",
                         Inches(0.8), Inches(1.7), Inches(4.2), Inches(5.4), radius=16)
    except Exception as e:
        print(f"  [WARN] 해상 이미지 실패: {e}")
    
    # 오버레이 (반투명 primary)
    overlay = add_rectangle(slide, Inches(0.8), Inches(1.7), Inches(4.2), Inches(5.4),
                            fill_color=COLORS["primary"], radius=16)
    overlay.fill.solid()
    overlay.fill.fore_color.rgb = hex_to_rgb(COLORS["primary"])
    from pptx.oxml.ns import qn
    spPr = overlay._sp.spPr
    solidFill = spPr.find(qn('a:solidFill'))
    if solidFill is not None:
        srgbClr = solidFill.find(qn('a:srgbClr'))
        if srgbClr is not None:
            from lxml import etree
            alpha = etree.SubElement(srgbClr, qn('a:alpha'))
            alpha.set('val', '75000')  # 75% 불투명
    
    add_text_box(slide, Inches(1.0), Inches(5.4), Inches(3.8), Inches(0.5),
                 "🚢", font_size=32)
    add_text_box(slide, Inches(1.0), Inches(5.95), Inches(3.8), Inches(0.5),
                 "해상 포워딩", font_size=26, font_color="FFFFFF", bold=True)
    add_text_box(slide, Inches(1.0), Inches(6.45), Inches(3.8), Inches(0.3),
                 "Ocean Freight Forwarding", font_size=12, font_color="FFFFFF")
    add_text_box(slide, Inches(1.0), Inches(6.75), Inches(3.8), Inches(0.3),
                 "FCL/LCL 수출입, 정기선 서비스", font_size=11, font_color="FFFFFF")
    
    # 중형 카드 1: 항공 포워딩 (우측 상단 좌)
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1436491865332-7a61a109cc05?w=400&h=300&fit=crop",
                         Inches(5.2), Inches(1.7), Inches(3.8), Inches(2.5), radius=16)
    except Exception as e:
        print(f"  [WARN] 항공 이미지 실패: {e}")
    
    overlay2 = add_rectangle(slide, Inches(5.2), Inches(1.7), Inches(3.8), Inches(2.5),
                             fill_color=COLORS["accent"], radius=16)
    overlay2.fill.solid()
    overlay2.fill.fore_color.rgb = hex_to_rgb(COLORS["accent"])
    spPr2 = overlay2._sp.spPr
    solidFill2 = spPr2.find(qn('a:solidFill'))
    if solidFill2 is not None:
        srgbClr2 = solidFill2.find(qn('a:srgbClr'))
        if srgbClr2 is not None:
            alpha2 = etree.SubElement(srgbClr2, qn('a:alpha'))
            alpha2.set('val', '75000')
    
    add_text_box(slide, Inches(5.4), Inches(2.9), Inches(3.4), Inches(0.4),
                 "✈️", font_size=28)
    add_text_box(slide, Inches(5.4), Inches(3.3), Inches(3.4), Inches(0.4),
                 "항공 포워딩", font_size=20, font_color="FFFFFF", bold=True)
    add_text_box(slide, Inches(5.4), Inches(3.7), Inches(3.4), Inches(0.25),
                 "Air Freight Forwarding", font_size=11, font_color="FFFFFF")
    
    # 중형 카드 2: 내륙 운송 (우측 상단 우)
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1519003722824-194d4455a60c?w=400&h=300&fit=crop",
                         Inches(9.2), Inches(1.7), Inches(3.8), Inches(2.5), radius=16)
    except Exception as e:
        print(f"  [WARN] 내륙 이미지 실패: {e}")
    
    overlay3 = add_rectangle(slide, Inches(9.2), Inches(1.7), Inches(3.8), Inches(2.5),
                             fill_color=COLORS["success"], radius=16)
    overlay3.fill.solid()
    overlay3.fill.fore_color.rgb = hex_to_rgb(COLORS["success"])
    spPr3 = overlay3._sp.spPr
    solidFill3 = spPr3.find(qn('a:solidFill'))
    if solidFill3 is not None:
        srgbClr3 = solidFill3.find(qn('a:srgbClr'))
        if srgbClr3 is not None:
            alpha3 = etree.SubElement(srgbClr3, qn('a:alpha'))
            alpha3.set('val', '75000')
    
    add_text_box(slide, Inches(9.4), Inches(2.9), Inches(3.4), Inches(0.4),
                 "🚛", font_size=28)
    add_text_box(slide, Inches(9.4), Inches(3.3), Inches(3.4), Inches(0.4),
                 "내륙 운송", font_size=20, font_color="FFFFFF", bold=True)
    add_text_box(slide, Inches(9.4), Inches(3.7), Inches(3.4), Inches(0.25),
                 "Inland Transportation", font_size=11, font_color="FFFFFF")
    
    # 중형 카드 3: 프로젝트 카고 (우측 하단, 가로로 길게)
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1504307651254-35680f356dfd?w=400&h=300&fit=crop",
                         Inches(5.2), Inches(4.4), Inches(7.8), Inches(2.7), radius=16)
    except Exception as e:
        print(f"  [WARN] 프로젝트 이미지 실패: {e}")
    
    overlay4 = add_rectangle(slide, Inches(5.2), Inches(4.4), Inches(7.8), Inches(2.7),
                             fill_color=COLORS["highlight"], radius=16)
    overlay4.fill.solid()
    overlay4.fill.fore_color.rgb = hex_to_rgb(COLORS["highlight"])
    spPr4 = overlay4._sp.spPr
    solidFill4 = spPr4.find(qn('a:solidFill'))
    if solidFill4 is not None:
        srgbClr4 = solidFill4.find(qn('a:srgbClr'))
        if srgbClr4 is not None:
            alpha4 = etree.SubElement(srgbClr4, qn('a:alpha'))
            alpha4.set('val', '75000')
    
    add_text_box(slide, Inches(5.4), Inches(5.7), Inches(7.4), Inches(0.4),
                 "📦", font_size=28)
    add_text_box(slide, Inches(5.4), Inches(6.1), Inches(7.4), Inches(0.4),
                 "프로젝트 카고", font_size=20, font_color="FFFFFF", bold=True)
    add_text_box(slide, Inches(5.4), Inches(6.5), Inches(7.4), Inches(0.25),
                 "Project & Heavy Cargo", font_size=11, font_color="FFFFFF")
    add_text_box(slide, Inches(5.4), Inches(6.75), Inches(7.4), Inches(0.25),
                 "중장비, 기계류, 위험물, 부정기선", font_size=11, font_color="FFFFFF")


def create_slide_6_support(prs, slide_data, design_tokens):
    """슬라이드 6: 부가서비스 (Split Layout) - 하드코딩"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["backgroundAlt"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
                 "SUPPORT INFRASTRUCTURE", font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목 (줄바꿈 포함)
    add_text_box(slide, Inches(0.8), Inches(0.95), Inches(6), Inches(1),
                 "파트너사를 위한\n물류 인프라 지원",
                 font_size=36, font_color=COLORS["text"], bold=True)
    
    # 설명
    add_text_box(slide, Inches(0.8), Inches(2.1), Inches(5.5), Inches(0.6),
                 "포워딩 파트너사의 물류망 확장을 지원하는 든든한 인프라가 되겠습니다.",
                 font_size=16, font_color=COLORS["textLight"])
    
    # 체크리스트 (4개 항목)
    checklist_items = [
        ("물류 거점 보관 지원", "일반/보세/냉장/냉동 창고 보관 서비스"),
        ("통관 및 서류 대행", "수출입 통관 원스톱 서비스"),
        ("배송 거점 제공", "협력사 임시 보관 및 배송 거점"),
        ("SCM/3PL 지원", "공급망 관리 파트너 서비스"),
    ]
    
    list_top = Inches(3.0)
    
    for i, (title, desc) in enumerate(checklist_items):
        y = list_top + i * Inches(0.95)
        
        # 체크 아이콘 배경
        add_rectangle(slide, Inches(0.8), y, Inches(0.45), Inches(0.45),
                      fill_color="EEF0FE", radius=50)
        add_text_box(slide, Inches(0.8), y, Inches(0.45), Inches(0.45),
                     "✓", font_size=16, font_color=COLORS["primary"], bold=True, align="center")
        
        # 제목
        add_text_box(slide, Inches(1.4), y, Inches(4.5), Inches(0.35),
                     title, font_size=18, font_color=COLORS["text"], bold=True)
        
        # 설명
        add_text_box(slide, Inches(1.4), y + Inches(0.35), Inches(4.5), Inches(0.35),
                     desc, font_size=16, font_color=COLORS["textLight"])
    
    # 우측 이미지
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1553413077-190dd305871c?w=600&h=400&fit=crop",
                         Inches(7.0), Inches(1.2), Inches(5.8), Inches(5.8), radius=16)
    except Exception as e:
        print(f"  [WARN] 창고 이미지 실패: {e}")


def create_slide_7_network(prs, slide_data, design_tokens):
    """슬라이드 7: 글로벌 네트워크 (Map + Stats) - 하드코딩"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["secondary"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.5), Inches(0.4), Inches(12.3), Inches(0.4),
                 "GLOBAL NETWORK", font_size=14, font_color=COLORS["primary"], bold=True, align="center")
    
    # 제목
    add_text_box(slide, Inches(0.5), Inches(0.9), Inches(12.3), Inches(0.6),
                 "전 세계를 잇는 물류 네트워크",
                 font_size=38, font_color="FFFFFF", bold=True, align="center")
    
    # 지도 이미지
    try:
        add_rounded_image(slide, "https://images.unsplash.com/photo-1526778548025-fa2f459cd5c1?w=1200&h=400&fit=crop",
                         Inches(1.5), Inches(1.7), Inches(10.3), Inches(3.5), radius=16)
    except Exception as e:
        print(f"  [WARN] 지도 이미지 실패: {e}")
    
    # 중앙 통계 오버레이 (30+ 파트너 국가)
    add_rectangle(slide, Inches(5.2), Inches(2.8), Inches(3), Inches(1.3),
                  fill_color="0F172A", radius=16)
    add_text_box(slide, Inches(5.2), Inches(2.95), Inches(3), Inches(0.7),
                 "30+",
                 font_size=40, font_color=COLORS["primary"], bold=True, align="center")
    add_text_box(slide, Inches(5.2), Inches(3.65), Inches(3), Inches(0.4),
                 "파트너 국가",
                 font_size=14, font_color=COLORS["textMuted"], align="center")
    
    # 지역 뱃지들 (5개)
    regions = [
        ("Asia", "8"),
        ("Europe", "8"),
        ("Middle East", "7"),
        ("Americas", "6"),
        ("Africa", "2"),
    ]
    
    badge_width = Inches(2.2)
    badge_height = Inches(1.0)
    start_x = Inches(1.0)
    badge_top = Inches(5.6)
    
    for i, (name, count) in enumerate(regions):
        x = start_x + i * (badge_width + Inches(0.2))
        
        add_rectangle(slide, x, badge_top, badge_width, badge_height,
                      fill_color="1E293B", radius=8)
        add_text_box(slide, x, badge_top + Inches(0.15), badge_width, Inches(0.35),
                     name,
                     font_size=14, font_color="FFFFFF", bold=True, align="center")
        add_text_box(slide, x, badge_top + Inches(0.45), badge_width, Inches(0.35),
                     f"{count}개국",
                     font_size=20, font_color=COLORS["primary"], bold=True, align="center")


def create_slide_8_partners(prs, slide_data, design_tokens):
    """슬라이드 8: 주요고객사"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["background"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(3), Inches(0.4),
                 "TRUSTED BY",
                 font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목
    add_text_box(slide, Inches(0.8), Inches(0.85), Inches(10), Inches(1.0),
                 "20개 이상의 기업이\n이음로지스틱과 함께합니다",
                 font_size=36, font_color=COLORS["text"], bold=True)
    
    # 설명
    add_text_box(slide, Inches(0.8), Inches(1.95), Inches(10), Inches(0.5),
                 "다양한 산업군의 주요 기업들에게 토탈 물류 서비스를 제공하고 있습니다.",
                 font_size=16, font_color=COLORS["textLight"])
    
    # 파트너 로고 그리드 (5x4)
    partners = [
        "가온전선", "롯데제과", "팔도", "HY", "삼양식품",
        "유앤아이원", "JK인터내셔널", "YC리싸이클", "에코비트프리텍", "롯데웰푸드",
        "CM코아", "광메탈", "페트로마인", "나우리소스", "삼성헬스토리",
        "대건비철", "팔인터내셔널", "신영금속", "디씨팩", "브링스글로벌",
    ]
    
    cols = 5
    rows = 4
    cell_width = Inches(2.3)
    cell_height = Inches(0.95)
    start_x = Inches(0.8)
    start_y = Inches(2.7)
    
    for i, partner in enumerate(partners):
        row = i // cols
        col = i % cols
        x = start_x + col * cell_width
        y = start_y + row * cell_height
        
        add_rectangle(slide, x, y, cell_width - Inches(0.15), cell_height - Inches(0.15),
                      fill_color=COLORS["backgroundAlt"], radius=8)
        add_text_box(slide, x, y + Inches(0.25), cell_width - Inches(0.15), Inches(0.5),
                     partner,
                     font_size=13, font_color=COLORS["text"], bold=True, align="center")


def create_slide_9_portfolio(prs, slide_data, design_tokens):
    """슬라이드 9: 프로젝트 포트폴리오 (3x2 Gallery) - 하드코딩"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["backgroundAlt"])
    
    # 섹션 라벨
    add_text_box(slide, Inches(0.8), Inches(0.4), Inches(5), Inches(0.4),
                 "PROJECT PORTFOLIO", font_size=14, font_color=COLORS["primary"], bold=True)
    
    # 제목
    add_text_box(slide, Inches(0.8), Inches(0.9), Inches(10), Inches(0.6),
                 "프로젝트 수행 실적", font_size=36, font_color=COLORS["text"], bold=True)
    
    # 부제목
    add_text_box(slide, Inches(0.8), Inches(1.55), Inches(10), Inches(0.5),
                 "중장비, 기계류, 특수화물 등 다양한 프로젝트 카고 경험",
                 font_size=16, font_color=COLORS["textLight"])
    
    # 프로젝트 카드 데이터 (6개)
    projects = [
        ("https://images.unsplash.com/photo-1494412574643-ff11b0a5c1c3?w=400&h=200&fit=crop", "해상운송"),
        ("https://images.unsplash.com/photo-1578575437130-527eed3abbec?w=400&h=200&fit=crop", "프로젝트 카고"),
        ("https://images.unsplash.com/photo-1504307651254-35680f356dfd?w=400&h=200&fit=crop", "항만 하역"),
        ("https://images.unsplash.com/photo-1519003722824-194d4455a60c?w=400&h=200&fit=crop", "내륙운송"),
        ("https://images.unsplash.com/photo-1601584115197-04ecc0da31d7?w=400&h=200&fit=crop", "특수화물"),
        ("https://images.unsplash.com/photo-1586528116311-ad8dd3c8310d?w=400&h=200&fit=crop", "특수화물"),
    ]
    
    card_width = Inches(3.9)
    card_height = Inches(2.5)
    gap = Inches(0.3)
    start_x = Inches(0.8)
    start_y = Inches(2.3)
    
    for i, (image_url, category) in enumerate(projects):
        row = i // 3
        col = i % 3
        x = start_x + col * (card_width + gap)
        y = start_y + row * (card_height + gap)
        
        # 카드 배경 (흰색)
        add_rectangle(slide, x, y, card_width, card_height,
                      fill_color="FFFFFF", radius=16)
        
        # 이미지 영역
        image_height = card_height - Inches(0.6)
        try:
            add_rounded_image(slide, image_url, x, y, card_width, image_height, radius=16)
        except Exception as e:
            print(f"  [WARN] 프로젝트 이미지 {i+1} 실패: {e}")
            # 실패 시 회색 배경
            add_rectangle(slide, x, y, card_width, image_height,
                          fill_color="E2E8F0", radius=16)
        
        # 하단 텍스트 영역 (흰색 배경)
        add_rectangle(slide, x, y + image_height, card_width, Inches(0.6),
                      fill_color="FFFFFF", radius=0)
        
        # 카테고리 라벨
        add_text_box(slide, x + Inches(0.2), y + card_height - Inches(0.45), 
                     card_width - Inches(0.4), Inches(0.35),
                     category, font_size=12, font_color=COLORS["primary"], bold=True)


def create_slide_10_message(prs, slide_data, design_tokens):
    """슬라이드 10: 회사 메시지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["primary"])
    
    # 인용문
    add_text_box(slide, Inches(1), Inches(2.5), Inches(11.3), Inches(1.5),
                 "우리는 크지 않습니다.\n하지만, 우리의 능력은 작지 않습니다.",
                 font_size=38, font_color="FFFFFF", bold=True, align="center")
    
    # 설명
    add_text_box(slide, Inches(1.5), Inches(4.3), Inches(10.3), Inches(2),
                 "저희는 경쟁이 아닌 협력을 추구합니다.\n포워딩 파트너사의 물류망 확장을 지원하는 든든한 인프라가 되겠습니다.\n무엇을 해야 할지, 어떻게 해야 할지, 어디로 가야 할지 잘 알고 있습니다.",
                 font_size=18, font_color="FFFFFF", align="center", line_spacing=1.8)


def create_slide_11_closing(prs, slide_data, design_tokens):
    """슬라이드 11: 감사 페이지"""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    set_slide_background(slide, COLORS["secondary"])
    
    # Thank You
    add_text_box(slide, Inches(0.5), Inches(2.2), Inches(12.3), Inches(1),
                 "Thank You",
                 font_size=56, font_color="FFFFFF", bold=True, align="center")
    
    # 태그라인
    add_text_box(slide, Inches(0.5), Inches(3.2), Inches(12.3), Inches(0.5),
                 "신뢰를 잇는 물류, 세계를 잇는 파트너",
                 font_size=18, font_color=COLORS["textMuted"], align="center")
    
    # 로고 (원형)
    add_rectangle(slide, Inches(6.0), Inches(4.0), Inches(1.3), Inches(1.3),
                  fill_color=COLORS["primary"], radius=50)
    add_text_box(slide, Inches(6.0), Inches(4.3), Inches(1.3), Inches(0.7),
                 "E", font_size=32, font_color="FFFFFF", bold=True, align="center")
    
    # 회사명
    add_text_box(slide, Inches(0.5), Inches(5.5), Inches(12.3), Inches(0.5),
                 "이음 로지스틱 주식회사",
                 font_size=22, font_color="FFFFFF", bold=True, align="center")
    
    # 연락처
    add_text_box(slide, Inches(0.5), Inches(6.2), Inches(12.3), Inches(0.4),
                 "🌐 www.e-umlk.com    ✉️ info@e-umlk.com    📞 +82 2 3662 8150",
                 font_size=14, font_color=COLORS["textMuted"], align="center")
    
    # Copyright
    add_text_box(slide, Inches(0.5), Inches(6.8), Inches(12.3), Inches(0.4),
                 "© 2025 EUM LOGISTIC CO., LTD. All rights reserved.",
                 font_size=12, font_color=COLORS["textLight"], align="center")


# ===== 메인 함수 =====

def convert_json_to_pptx(json_path: str, output_path: str = None):
    """JSON 파일을 편집 가능한 PPTX로 변환"""
    
    # 경로 처리
    json_path = Path(json_path)
    if not json_path.exists():
        print(f"[ERROR] JSON 파일을 찾을 수 없습니다: {json_path}")
        return False
    
    if output_path is None:
        output_path = json_path.parent / (json_path.stem + "_editable.pptx")
    else:
        output_path = Path(output_path)
    
    print(f"[Input] {json_path}")
    print(f"[Output] {output_path}")
    
    # JSON 로드
    with open(json_path, 'r', encoding='utf-8') as f:
        data = json.load(f)
    
    design_tokens = data.get("designTokens", {})
    slides_data = data.get("slides", [])
    
    print(f"[Found] {len(slides_data)} slides")
    
    # 프레젠테이션 생성
    prs = Presentation()
    prs.slide_width = SLIDE_WIDTH
    prs.slide_height = SLIDE_HEIGHT
    
    # 슬라이드 생성 함수 매핑
    slide_creators = {
        1: create_slide_1_cover,
        2: create_slide_2_company,
        3: create_slide_3_history,
        4: create_slide_4_values,
        5: create_slide_5_services,
        6: create_slide_6_support,
        7: create_slide_7_network,
        8: create_slide_8_partners,
        9: create_slide_9_portfolio,
        10: create_slide_10_message,
        11: create_slide_11_closing,
    }
    
    # 각 슬라이드 생성
    print("\n[Creating slides...]")
    slides_created = 0
    for slide_data in slides_data:
        slide_num = slide_data.get("slideNumber", 0)
        creator = slide_creators.get(slide_num)
        
        if creator:
            try:
                print(f"  [Slide {slide_num}] {slide_data.get('title', '')}")
                creator(prs, slide_data, design_tokens)
                slides_created += 1
            except Exception as e:
                print(f"  [ERROR] Slide {slide_num} 생성 실패: {e}")
                import traceback
                traceback.print_exc()
        else:
            # 범용 슬라이드 생성
            try:
                print(f"  [Slide {slide_num}] {slide_data.get('title', '')} (범용)")
                create_generic_slide(prs, slide_data, design_tokens)
                slides_created += 1
            except Exception as e:
                print(f"  [ERROR] Slide {slide_num} 범용 생성 실패: {e}")
    
    # 저장
    prs.save(str(output_path))
    print(f"\n[OK] 편집 가능한 PPTX 생성 완료: {output_path}")
    print(f"     Total {len(prs.slides)} slides")
    
    # 사용된 폰트를 프로젝트 폴더에 복사
    project_dir = json_path.parent
    copy_fonts_to_project(project_dir)
    
    return True


if __name__ == "__main__":
    if len(sys.argv) < 2:
        print("Usage: python json_to_pptx.py <presentation.json> [output.pptx]")
        print("Example: python json_to_pptx.py projects/eumlogistic/presentation.json")
        sys.exit(1)
    
    json_file = sys.argv[1]
    output_file = sys.argv[2] if len(sys.argv) > 2 else None
    
    success = convert_json_to_pptx(json_file, output_file)
    sys.exit(0 if success else 1)
